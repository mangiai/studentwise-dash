from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colours
NAVY    = RGBColor(0x1E, 0x3A, 0x5F)   # deep navy
TEAL    = RGBColor(0x0D, 0xB8, 0xA0)   # accent teal
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF0, 0xF4, 0xF8)   # slide background
DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # body text
GRAY    = RGBColor(0x64, 0x74, 0x8B)   # muted text
ORANGE  = RGBColor(0xF5, 0x9E, 0x0B)   # highlight

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=NAVY, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_bullet_box(slide, items, l, t, w, h, size=15, title=None,
                   bullet_color=TEAL, text_color=DARK):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 1)
        run.font.bold = True
        run.font.color.rgb = NAVY
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  ›  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = text_color

def slide_bg(slide, color=LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── SLIDE 1 — Title ──────────────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, NAVY)

# big teal accent bar left
add_rect(s, 0, 0, 0.35, 7.5, TEAL)

# product name
add_text(s, "UniFlow", 0.6, 1.6, 8, 1.4, size=72, bold=True, color=WHITE)
add_text(s, "University Portal", 0.6, 3.1, 8, 0.7, size=30, bold=False, color=TEAL)
add_text(s, "Student · Teacher · Admin · Moderator", 0.6, 3.85, 9, 0.5,
         size=17, color=RGBColor(0xB0, 0xC4, 0xDE), italic=True)

# tagline box bottom-right
add_rect(s, 9.0, 5.8, 4.0, 1.3, TEAL)
add_text(s, "Full-stack academic management\nReact 19 · Supabase · Vercel",
         9.15, 5.9, 3.7, 1.1, size=14, color=WHITE, align=PP_ALIGN.CENTER)


# ── SLIDE 2 — Problem & Solution ────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "The Problem  →  The Solution", 0.4, 0.15, 12, 0.8,
         size=28, bold=True, color=WHITE)

# Problem column
add_rect(s, 0.4, 1.3, 5.8, 5.6, WHITE)
add_rect(s, 0.4, 1.3, 5.8, 0.55, RGBColor(0xE5, 0x4D, 0x4D))
add_text(s, "Problem", 0.55, 1.33, 5.5, 0.5, size=17, bold=True, color=WHITE)
problems = [
    "Universities manage students, courses,\nfees & attendance across disconnected systems",
    "No single portal for all stakeholders",
    "Manual fee & challan processes",
    "No real-time attendance visibility",
    "Staff lack role-scoped access controls",
]
add_bullet_box(s, problems, 0.55, 2.0, 5.6, 4.5, size=14, text_color=DARK)

# Solution column
add_rect(s, 7.0, 1.3, 5.8, 5.6, WHITE)
add_rect(s, 7.0, 1.3, 5.8, 0.55, TEAL)
add_text(s, "Solution", 7.15, 1.33, 5.5, 0.5, size=17, bold=True, color=WHITE)
solutions = [
    "One portal: student, teacher, admin & moderator",
    "Role-based dashboards with live Supabase data",
    "Automated fee challan generation (PKR)",
    "Session-level attendance with real-time sync",
    "Granular RLS — admins delete, moderators edit",
]
add_bullet_box(s, solutions, 7.15, 2.0, 5.6, 4.5, size=14, text_color=DARK)


# ── SLIDE 3 — Key Features ──────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "Key Features", 0.4, 0.15, 12, 0.8, size=28, bold=True, color=WHITE)

features = [
    ("Dashboard & Analytics",    "GPA tracker, attendance ring charts,\nRecharts visualisations"),
    ("Courses & Enrolment",      "Browse enrolled courses, credit tracking,\nper-course attendance badges"),
    ("Attendance Tracking",      "Session-level marks, auto-synced summaries,\nshort-attendance alerts"),
    ("Fee Management",           "PKR semester balance, payment history,\ndownloadable fee challan"),
    ("Academic Results",         "Semester grades, CGPA, grade badges"),
    ("Staff Portal (Admin/Mod)", "CRUD students, teachers, courses;\nFee & notification management"),
    ("Real-time Sync",           "Supabase Realtime invalidates enrollments,\nfees, grades & attendance live"),
    ("Role-based Auth",          "4 roles · cookie sessions · RLS on every table"),
]

cols = 4
per_col = 2
box_w, box_h = 2.9, 2.3
gap_x, gap_y = 0.22, 0.25
start_x, start_y = 0.25, 1.2

for i, (title, body) in enumerate(features):
    col = i % cols
    row = i // cols
    x = start_x + col * (box_w + gap_x)
    y = start_y + row * (box_h + gap_y)
    add_rect(s, x, y, box_w, box_h, WHITE)
    add_rect(s, x, y, box_w, 0.42, TEAL)
    add_text(s, title, x + 0.1, y + 0.05, box_w - 0.15, 0.38,
             size=13, bold=True, color=WHITE)
    add_text(s, body, x + 0.1, y + 0.5, box_w - 0.2, box_h - 0.6,
             size=12, color=DARK, wrap=True)


# ── SLIDE 4 — User Roles ────────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "4 User Roles", 0.4, 0.15, 12, 0.8, size=28, bold=True, color=WHITE)

roles = [
    ("Student",    TEAL,   "/login",       [
        "Dashboard with GPA & attendance overview",
        "My Courses with attendance badges",
        "Fee management & challan download",
        "Academic results & CGPA",
        "Notifications & settings",
    ]),
    ("Teacher",    RGBColor(0x6D, 0x28, 0xD9), "/login", [
        "All student pages",
        "Reports & Analytics page",
        "Enrollment trends, fee status charts",
        "Department performance dashboards",
    ]),
    ("Admin",      RGBColor(0xE5, 0x4D, 0x4D), "/admin/login", [
        "Full staff portal (8 sections)",
        "Create / edit / DELETE students & teachers",
        "Mark attendance per session",
        "Manage fee challans, send notifications",
    ]),
    ("Moderator",  ORANGE, "/admin/login", [
        "Same staff portal as admin",
        "Create & edit — no delete permissions",
        "Hidden delete buttons enforced by RLS",
    ]),
]

box_w = 2.95
for i, (role, color, url, bullets) in enumerate(roles):
    x = 0.25 + i * (box_w + 0.2)
    add_rect(s, x, 1.2, box_w, 5.8, WHITE)
    add_rect(s, x, 1.2, box_w, 0.7, color)
    add_text(s, role, x + 0.1, 1.23, box_w - 0.2, 0.5, size=18, bold=True, color=WHITE)
    add_text(s, url, x + 0.1, 1.95, box_w - 0.2, 0.35, size=11,
             color=GRAY, italic=True)
    add_bullet_box(s, bullets, x + 0.1, 2.35, box_w - 0.2, 4.4, size=12,
                   text_color=DARK)


# ── SLIDE 5 — Student Portal Pages ─────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "Student Portal  —  All Pages", 0.4, 0.15, 12, 0.8,
         size=28, bold=True, color=WHITE)

pages = [
    ("/dashboard",      "Dashboard",         "GPA, attendance ring, course cards, fee alert"),
    ("/courses",        "My Courses",        "6 enrolled courses, credits, attendance badge"),
    ("/attendance",     "Attendance",        "Per-course % rings, short-attendance warning ❗"),
    ("/fees",           "Fee Management",    "PKR balance, payment history, Download Challan"),
    ("/teachers",       "Faculty Directory", "Teacher cards from live database"),
    ("/results",        "Academic Results",  "Semester grades, CGPA, grade colour badges"),
    ("/reports",        "Reports",           "Enrollment & fee charts (teacher / admin only)"),
    ("/notifications",  "Notifications",     "Fee, attendance, course, announcement alerts"),
    ("/settings",       "Settings",          "Profile, notification prefs, password change"),
]

for i, (route, title, desc) in enumerate(pages):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 4.35
    y = 1.25 + row * 1.95
    add_rect(s, x, y, 4.1, 1.75, WHITE)
    add_rect(s, x, y, 0.06, 1.75, TEAL)
    add_text(s, title, x + 0.2, y + 0.1, 3.8, 0.45, size=14, bold=True, color=NAVY)
    add_text(s, route, x + 0.2, y + 0.52, 3.8, 0.3, size=10, color=GRAY, italic=True)
    add_text(s, desc,  x + 0.2, y + 0.82, 3.8, 0.8, size=11, color=DARK, wrap=True)


# ── SLIDE 6 — Staff Portal ──────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "Staff Portal  —  Admin & Moderator", 0.4, 0.15, 12, 0.8,
         size=28, bold=True, color=WHITE)

sections = [
    ("/admin/dashboard",    "Overview",          "Summary stats & quick links"),
    ("/admin/students",     "Students",          "Create / edit / delete student records"),
    ("/admin/teachers",     "Teachers",          "Manage faculty — assign to courses"),
    ("/admin/courses",      "Courses",           "Create courses, assign instructors"),
    ("/admin/enrollments",  "Enrolments",        "Enrol / unenrol students in courses"),
    ("/admin/attendance",   "Attendance",        "Mark per-session attendance (Fall 2026)"),
    ("/admin/fees",         "Fees & Challans",   "Regenerate challans, set Paid/Pending/Overdue"),
    ("/admin/notifications","Notifications",     "Broadcast alerts to individual users"),
]

for i, (route, title, desc) in enumerate(sections):
    col = i % 4
    row = i // 4
    x = 0.3 + col * 3.2
    y = 1.25 + row * 2.55
    add_rect(s, x, y, 3.0, 2.3, WHITE)
    add_rect(s, x, y, 3.0, 0.5, RGBColor(0x1E, 0x3A, 0x5F))
    add_text(s, title, x + 0.1, y + 0.07, 2.8, 0.42, size=14, bold=True, color=WHITE)
    add_text(s, route, x + 0.1, y + 0.56, 2.8, 0.3, size=10, color=GRAY, italic=True)
    add_text(s, desc,  x + 0.1, y + 0.88, 2.8, 1.3, size=12, color=DARK, wrap=True)

# Moderator note
add_rect(s, 0.3, 6.05, 12.7, 0.95, RGBColor(0xFF, 0xF3, 0xCD))
add_text(s,
         "Moderator:  same 8 sections — create & edit allowed  |  Delete buttons hidden (enforced by RLS + useStaffPermissions())",
         0.45, 6.12, 12.4, 0.75, size=13, color=RGBColor(0x92, 0x40, 0x0E), bold=False)


# ── SLIDE 7 — Tech Stack ────────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "Technology Stack", 0.4, 0.15, 12, 0.8, size=28, bold=True, color=WHITE)

cols_data = [
    ("Frontend", TEAL, [
        "React 19.2",
        "TanStack Router 1.168",
        "TanStack Start (SSR)",
        "TanStack Query 5.83",
        "Tailwind CSS v4",
        "shadcn/ui + Radix UI",
        "Recharts 2.15",
        "React Hook Form + Zod",
    ]),
    ("Backend & DB", RGBColor(0x6D, 0x28, 0xD9), [
        "Supabase Postgres",
        "Supabase Auth (GoTrue)",
        "@supabase/ssr (cookies)",
        "Row Level Security (RLS)",
        "Supabase Realtime",
        "13 tables · 5 migrations",
        "pgcrypto (seed hashing)",
        "33 server functions",
    ]),
    ("Build & Deploy", RGBColor(0xE5, 0x4D, 0x4D), [
        "Vite 7.3 (bundler)",
        "Nitro 3.0 (serverless preset)",
        "Vercel (production host)",
        "Node.js 22.x runtime",
        "TypeScript 5.8",
        "ESLint + Prettier",
        "Supabase CLI (migrations)",
        "GitHub Actions (CI migrate)",
    ]),
]

bw = 3.8
for i, (label, color, items) in enumerate(cols_data):
    x = 0.4 + i * (bw + 0.45)
    add_rect(s, x, 1.2, bw, 5.8, WHITE)
    add_rect(s, x, 1.2, bw, 0.6, color)
    add_text(s, label, x + 0.15, 1.25, bw - 0.2, 0.52, size=17, bold=True, color=WHITE)
    add_bullet_box(s, items, x + 0.15, 1.9, bw - 0.2, 5.0, size=13, text_color=DARK)


# ── SLIDE 8 — Architecture ──────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "Architecture", 0.4, 0.15, 12, 0.8, size=28, bold=True, color=WHITE)

layers = [
    (TEAL,   "Browser  (React 19)",
     "TanStack Router · shadcn/ui · Recharts · Tailwind CSS v4\nRealtime hooks · MobileNavBar · NotificationBell"),
    (RGBColor(0x6D,0x28,0xD9), "Vercel  (Nitro Serverless Functions)",
     "TanStack Start server functions (SSR)\nAuth guards · data loaders · staff CRUD mutations"),
    (NAVY,   "Supabase Cloud",
     "Auth (GoTrue) · Postgres + RLS · Realtime subscriptions"),
]

arrow_x = 6.4
for i, (color, title, body) in enumerate(layers):
    y = 1.3 + i * 1.85
    add_rect(s, 0.5, y, 12.3, 1.55, WHITE)
    add_rect(s, 0.5, y, 0.25, 1.55, color)
    add_text(s, title, 0.9, y + 0.08, 11.5, 0.5, size=16, bold=True, color=color)
    add_text(s, body,  0.9, y + 0.58, 11.5, 0.85, size=13, color=DARK, wrap=True)
    if i < 2:
        add_text(s, "▼  HTTP / Cookies + Supabase JS Client",
                 3.5, y + 1.62, 6, 0.3, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# RLS note
add_rect(s, 0.5, 6.75, 12.3, 0.55, RGBColor(0xDC, 0xFD, 0xF0))
add_text(s, "Security:  RLS on all 13 tables · students read own rows only · staff via is_staff() · admins via is_admin()",
         0.7, 6.78, 12.0, 0.5, size=12, color=RGBColor(0x06, 0x53, 0x36))


# ── SLIDE 9 — Database Schema ───────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "Database  —  13 Tables · 5 Migrations", 0.4, 0.15, 12, 0.8,
         size=28, bold=True, color=WHITE)

tables = [
    ("departments",         "CS, EE, BBA, MATH"),
    ("profiles",            "User profiles + notification prefs (linked to auth.users)"),
    ("students",            "Records: ID, GPA, credits, fee status"),
    ("teachers",            "Faculty records"),
    ("courses",             "Catalog with instructor assignment"),
    ("enrollments",         "Student ↔ course per semester"),
    ("attendance_records",  "Auto-calculated % per enrolment (trigger-synced)"),
    ("class_sessions",      "Individual class dates (Fall 2026 term)"),
    ("session_attendance",  "Per-session present/absent marks"),
    ("fee_transactions",    "Payment history (PKR)"),
    ("semester_fees",       "Current semester fee balances"),
    ("course_grades",       "Semester grades & grade points"),
    ("notifications",       "User-specific alerts (fee, attendance, course, announcement)"),
]

for i, (name, desc) in enumerate(tables):
    col = i % 2
    row = i // 2
    x = 0.35 + col * 6.4
    y = 1.25 + row * 0.75
    add_rect(s, x, y, 6.15, 0.65, WHITE)
    add_rect(s, x, y, 1.85, 0.65, RGBColor(0xE0, 0xF2, 0xFE))
    add_text(s, name, x + 0.08, y + 0.1, 1.7, 0.45, size=12, bold=True,
             color=NAVY)
    add_text(s, desc, x + 2.0, y + 0.1, 4.0, 0.45, size=11, color=DARK)

# realtime note
add_rect(s, 0.35, 7.05, 12.6, 0.38, RGBColor(0xDC, 0xFD, 0xF0))
add_text(s,
         "Realtime publications:  enrollments · notifications · semester_fees · course_grades · attendance_records · class_sessions · session_attendance",
         0.5, 7.07, 12.3, 0.32, size=10, color=RGBColor(0x06, 0x53, 0x36))


# ── SLIDE 10 — Demo Flow ────────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "Demo Walkthrough", 0.4, 0.15, 12, 0.8, size=28, bold=True, color=WHITE)

steps = [
    ("1", TEAL,                         "Student Login",
     "sarah@studentwise.test / StudentWise123!\n→  /login"),
    ("2", RGBColor(0x6D,0x28,0xD9),    "Dashboard",
     "GPA, attendance ring charts, course progress (live)"),
    ("3", TEAL,                         "Attendance",
     "Ring charts · short-attendance warning on OS (68%)"),
    ("4", RGBColor(0x6D,0x28,0xD9),    "Fees",
     "PKR balance · payment history · Download Challan"),
    ("5", TEAL,                         "Results + Notifications",
     "Semester grades, CGPA · 2 unread alerts"),
    ("6", RGBColor(0xE5,0x4D,0x4D),    "Teacher View",
     "teacher@studentwise.test  →  Reports & Analytics"),
    ("7", RGBColor(0xE5,0x4D,0x4D),    "Admin Portal",
     "admin@studentwise.test  →  /admin/login\nCRUD student, mark attendance, update fee status"),
    ("8", ORANGE,                       "Moderator",
     "moderator@studentwise.test\nSame portal — delete buttons hidden"),
]

bw, bh = 2.85, 2.6
for i, (num, color, title, desc) in enumerate(steps):
    col = i % 4
    row = i // 4
    x = 0.3 + col * (bw + 0.3)
    y = 1.25 + row * (bh + 0.3)
    add_rect(s, x, y, bw, bh, WHITE)
    add_rect(s, x, y, 0.45, bh, color)
    add_text(s, num, x + 0.06, y + bh/2 - 0.28, 0.35, 0.55,
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x + 0.6, y + 0.15, bw - 0.7, 0.5,
             size=14, bold=True, color=NAVY)
    add_text(s, desc,  x + 0.6, y + 0.65, bw - 0.7, bh - 0.75,
             size=11, color=DARK, wrap=True)


# ── SLIDE 11 — Roadmap ──────────────────────────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, LIGHT)
add_rect(s, 0, 0, 13.33, 1.1, NAVY)
add_text(s, "Roadmap  —  Future Enhancements", 0.4, 0.15, 12, 0.8,
         size=28, bold=True, color=WHITE)

roadmap = [
    ("High",   RGBColor(0xE5,0x4D,0x4D), "Fee Payment Gateway",
     "Integrate online payment for semester fee challans"),
    ("High",   RGBColor(0xE5,0x4D,0x4D), "PDF Export",
     "Downloadable PDF for results & fee receipts (challan UI ready)"),
    ("Medium", ORANGE,                    "Email Delivery",
     "Wire actual email sending to the notification preference toggles"),
    ("Medium", ORANGE,                    "Teacher Self-Service Attendance",
     "Let teachers mark their own courses instead of staff-only"),
    ("Medium", ORANGE,                    "Bulk Attendance Import",
     "CSV upload for session attendance records"),
    ("Low",    TEAL,                      "Push Notifications",
     "Browser / mobile push (preference toggle already saved)"),
    ("Low",    TEAL,                      "Multi-semester Attendance",
     "Expand beyond Fall 2026 term to full attendance history"),
    ("Low",    TEAL,                      "Student Self-Registration",
     "Allow students to link their own ID at signup without admin approval"),
]

for i, (priority, color, title, desc) in enumerate(roadmap):
    col = i % 2
    row = i // 2
    x = 0.35 + col * 6.4
    y = 1.25 + row * 1.45
    add_rect(s, x, y, 6.15, 1.3, WHITE)
    add_rect(s, x, y, 0.25, 1.3, color)
    add_rect(s, x + 0.35, y + 0.1, 1.0, 0.38, color)
    add_text(s, priority, x + 0.38, y + 0.12, 0.95, 0.35,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x + 1.5, y + 0.08, 4.5, 0.42, size=14, bold=True, color=NAVY)
    add_text(s, desc,  x + 1.5, y + 0.52, 4.5, 0.65, size=12, color=DARK, wrap=True)


# ── SLIDE 12 — Closing / Quick Reference ───────────────────────────────────

s = prs.slides.add_slide(BLANK)
slide_bg(s, NAVY)
add_rect(s, 0, 0, 0.35, 7.5, TEAL)

add_text(s, "UniFlow", 0.6, 0.5, 8, 1.0, size=54, bold=True, color=WHITE)
add_text(s, "University Portal — Quick Reference", 0.6, 1.55, 10, 0.55,
         size=20, color=TEAL, bold=True)

# quick ref box
add_rect(s, 0.6, 2.2, 12.3, 4.6, RGBColor(0x0F, 0x23, 0x44))

creds = [
    ("Student",   "sarah@studentwise.test",      "/login"),
    ("Teacher",   "teacher@studentwise.test",    "/login"),
    ("Admin",     "admin@studentwise.test",      "/admin/login"),
    ("Moderator", "moderator@studentwise.test",  "/admin/login"),
]

add_text(s, "Test Credentials  (password: StudentWise123!)",
         0.85, 2.35, 11.5, 0.45, size=14, bold=True, color=TEAL)

for i, (role, email, url) in enumerate(creds):
    y = 2.9 + i * 0.6
    add_text(s, f"{role:<12}", 0.9, y, 2.0, 0.5, size=13, color=ORANGE, bold=True)
    add_text(s, email,         2.9, y, 5.5, 0.5, size=13, color=WHITE)
    add_text(s, url,           8.4, y, 4.0, 0.5, size=13, color=TEAL, italic=True)

add_text(s, "Stack:  React 19 · TanStack Start · Supabase Postgres + Auth + Realtime · Vercel · TypeScript",
         0.85, 5.55, 11.5, 0.45, size=12, color=RGBColor(0xB0, 0xC4, 0xDE))
add_text(s, "Repo:   github.com/mangiai/studentwise-dash",
         0.85, 6.05, 11.5, 0.4, size=12, color=RGBColor(0xB0, 0xC4, 0xDE))


# ── save ────────────────────────────────────────────────────────────────────

out = "/Users/mangiai/Downloads/studentwise-dash-main/UniFlow_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
